import streamlit as st
import pandas as pd
from jinja2 import Environment, FileSystemLoader
from weasyprint import HTML, CSS
from datetime import datetime
import streamlit.components.v1 as components
from openai import OpenAI
import json
import os
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import io
import re
import PyPDF2
from pptx import Presentation
import openpyxl

# --- 학습된 문서 관리 ---
learned_documents = {}
learning_status = {"manual": False, "samples": False}

def load_learned_documents():
    """학습된 문서 내용을 로드합니다."""
    global learned_documents, learning_status
    try:
        if os.path.exists('learned_documents.json'):
            with open('learned_documents.json', 'r', encoding='utf-8') as f:
                learned_documents = json.load(f)
                
                # 기존 방식과 새로운 방식 모두 지원
                learning_status = {
                    "manual": learned_documents.get('manual', {}).get('content', '') != '',
                    "samples": learned_documents.get('samples', {}).get('content', '') != ''
                }
                
                # 새로운 files 구조가 있으면 추가로 확인
                if learned_documents.get('files'):
                    files_data = learned_documents.get('files', {})
                    successful_files = [f for f, data in files_data.items() if data.get('success')]
                    if successful_files:
                        learning_status["files_learned"] = True
                    else:
                        learning_status["files_learned"] = False
                
                return True
    except Exception as e:
        st.error(f"학습된 문서를 로드하는 중 오류가 발생했습니다: {str(e)}")
    return False

def get_learning_enhanced_prompt(base_prompt, doc_type):
    """학습된 내용이 포함된 강화된 프롬프트를 생성합니다."""
    if not learned_documents:
        return base_prompt
    
    enhancement = "\n\n[학습된 문서 가이드라인]:\n"
    total_content = ""
    
    # 기존 manual, samples 키 지원
    if learning_status.get("manual") and learned_documents.get('manual', {}).get('content'):
        enhancement += "\n📋 문서작성 가이드라인:\n"
        enhancement += learned_documents['manual']['content'][:2000]  # 2000자로 확장
    
    if learning_status.get("samples") and learned_documents.get('samples', {}).get('content'):
        enhancement += "\n📝 품의서 작성 패턴:\n"
        enhancement += learned_documents['samples']['content'][:2000]  # 2000자로 확장
    
    # 새로운 files 구조 지원 - 문서 유형별로 관련성 높은 파일 우선 포함
    if learned_documents.get('files'):
        relevant_files = []
        other_files = []
        
        for filename, file_data in learned_documents['files'].items():
            if file_data.get('success') and file_data.get('content'):
                # 현재 작성 중인 문서 유형과 관련성 체크
                is_relevant = False
                if doc_type == '품의서' and ('품의서' in filename or '모음' in filename or '메뉴얼' in filename):
                    is_relevant = True
                elif doc_type == '공지문' and ('공지' in filename or '메뉴얼' in filename):
                    is_relevant = True
                elif doc_type == '공문' and ('공문' in filename or '메뉴얼' in filename):
                    is_relevant = True
                elif doc_type == '비즈니스 이메일' and ('이메일' in filename or 'email' in filename.lower() or '메뉴얼' in filename):
                    is_relevant = True
                
                if is_relevant:
                    relevant_files.append((filename, file_data))
                else:
                    other_files.append((filename, file_data))
        
        # 관련 파일을 먼저 포함
        all_files = relevant_files + other_files
        
        if all_files:
            enhancement += "\n📚 학습된 전문 문서 가이드라인:\n"
            
            for filename, file_data in all_files[:5]:  # 최대 5개 파일만 포함
                # 파일명에서 카테고리 추론
                if '메뉴얼' in filename or 'manual' in filename.lower():
                    category = "📋 작성 가이드라인"
                elif '품의서' in filename or '모음' in filename:
                    category = "📝 품의서 실제 사례"
                elif '공지' in filename:
                    category = "📢 공지문 템플릿"
                elif '공문' in filename:
                    category = "📄 공문 양식"
                elif '이메일' in filename or 'email' in filename.lower():
                    category = "📧 이메일 양식"
                else:
                    category = "📖 참고 문서"
                
                enhancement += f"\n{category}:\n"
                
                # 내용을 더 길게 포함 (문서 유형 관련성에 따라 조정)
                content = file_data['content']
                if filename in [f[0] for f in relevant_files]:
                    max_length = 3000  # 관련성 높은 파일은 더 길게
                else:
                    max_length = 1500  # 일반 파일은 중간 길이
                
                if len(content) > max_length:
                    # 중요한 부분을 보존하기 위해 앞부분과 뒷부분을 포함
                    front_part = content[:max_length//2]
                    back_part = content[-(max_length//2):]
                    content = front_part + "\n...(중간 내용 생략)...\n" + back_part
                
                enhancement += content + "\n"
    
    enhancement += f"\n\n위의 모든 학습된 가이드라인과 실제 사례를 바탕으로 '{doc_type}' 문서의 전문성과 완성도를 최대한 높여 작성해주세요. 특히 학습된 문서의 구조, 문체, 표현 방식을 참고하여 한국 비즈니스 문서 표준에 맞춰 작성하세요."
    
    return base_prompt + enhancement

def reset_learning_data():
    """학습 데이터를 초기화합니다."""
    global learned_documents, learning_status
    try:
        if os.path.exists('learned_documents.json'):
            os.remove('learned_documents.json')
        learned_documents = {}
        learning_status = {"manual": False, "samples": False}
        return True
    except Exception as e:
        st.sidebar.error(f"❌ 초기화 중 오류: {str(e)}")
        return False

# 앱 시작 시 학습된 문서 로드
load_learned_documents()

# --- 모델 설정 관리 ---
if 'selected_model' not in st.session_state:
    st.session_state.selected_model = "gpt-4o-mini"

if 'model_password_verified' not in st.session_state:
    st.session_state.model_password_verified = False

# --- AI 설정 ---
client = None
openai_available = False

try:
    if "OPENAI_API_KEY" in st.secrets:
        client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])
        openai_available = True
    else:
        st.warning("⚠️ OpenAI API 키가 설정되지 않았습니다. AI 기능이 비활성화됩니다.")
except Exception as e:
    st.error(f"OpenAI 클라이언트 초기화 중 오류가 발생했습니다: {str(e)}")
    st.warning("AI 기능이 비활성화됩니다.")

def get_ai_response(system_prompt, user_prompt):
    """OpenAI API를 호출하는 범용 함수"""
    if not openai_available or client is None:
        st.error("⚠️ OpenAI API가 설정되지 않아 AI 기능을 사용할 수 없습니다.")
        return None
        
    if not system_prompt or not user_prompt:
        st.error("프롬프트가 비어있습니다.")
        return None
        
    try:
        response = client.chat.completions.create(
            model=st.session_state.selected_model,
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7,
            max_tokens=3000,
            timeout=30
        )
        
        if not response.choices or not response.choices[0].message.content:
            st.error("AI로부터 응답을 받지 못했습니다.")
            return None
            
        content = response.choices[0].message.content.strip()
        if not content:
            st.error("AI 응답이 비어있습니다.")
            return None
            
        return json.loads(content)
        
    except json.JSONDecodeError as e:
        st.error(f"AI 응답 형식이 올바르지 않습니다: {str(e)}")
        return None
    except Exception as e:
        error_msg = str(e)
        if "rate limit" in error_msg.lower():
            st.error("⚠️ API 요청 한도를 초과했습니다. 잠시 후 다시 시도해주세요.")
        elif "timeout" in error_msg.lower():
            st.error("⚠️ AI 응답 시간이 초과되었습니다. 다시 시도해주세요.")
        elif "insufficient_quota" in error_msg.lower():
            st.error("⚠️ OpenAI API 할당량이 부족합니다. 계정을 확인해주세요.")
        else:
            st.error(f"AI 생성 중 오류가 발생했습니다: {error_msg}")
        return None

def analyze_keywords(keywords, doc_type):
    """키워드를 분석하여 추가 질문을 생성하는 함수"""
    analysis_prompt = f"사용자가 '{doc_type}' 작성을 위해 다음 키워드를 입력했습니다: '{keywords}'. 6W3H 원칙에 따라 완성도 높은 문서를 작성하기에 정보가 부족하다면, 가장 중요한 질문 2-3개를 `{{\"status\": \"incomplete\", \"questions\": [\"질문1\", \"질문2\"]}}` 형식으로 반환하고, 충분하다면 `{{\"status\": \"complete\"}}` 를 반환하세요."
    base_system_prompt = "당신은 사용자의 입력을 분석하여 문서 작성에 필요한 추가 정보를 질문하는 시스템입니다. 반드시 지정된 JSON 형식으로만 응답해야 합니다."
    
    # 학습된 내용으로 시스템 프롬프트 강화
    enhanced_system_prompt = get_learning_enhanced_prompt(base_system_prompt, doc_type)
    
    return get_ai_response(enhanced_system_prompt, analysis_prompt)

def generate_ai_draft(doc_type, context_keywords, file_context=""):
    """최종 키워드와 파일 내용을 바탕으로 AI 초안을 생성하는 함수"""
    user_prompt = f"다음 정보를 바탕으로 '{doc_type}' 초안을 JSON 형식으로 생성해주세요:\n\n[핵심 키워드]: {context_keywords}\n\n[첨부 파일 내용]:\n{file_context}"
    # 기본 프롬프트를 학습된 내용으로 강화
    base_prompts = {
        "품의서": "당신은 한국의 '주식회사 몬쉘코리아' 소속의 유능한 사원입니다. 품의서 초안을 생성합니다. **절대 규칙**: 'body' 필드는 반드시 다음 형식으로 작성하세요:\n\n1. 첫 번째 주요 항목\n  1) 세부 사항\n    (1) 구체적 내용\n  2) 추가 세부 사항\n2. 두 번째 주요 항목\n  1) 세부 사항\n\n이런 식으로 `1.`, `  1)`, `    (1)` 구분기호를 의무적으로 사용하여 체계적으로 작성하세요. 절대로 구분기호 없이 단순 문장 나열하지 마세요. 문장 종결어미는 `...함.`, `...요청함.` 형태로 하고, 각 문장 마침표 후 줄바꿈하세요. \n\n**중요**: 'items' 필드는 사용자가 제공한 키워드에서 표로 정리할 수 있는 구체적인 데이터가 있을 때만 생성하세요. 예를 들어:\n- 구매 품목: 항목명, 수량, 단가, 총액 등\n- 직급별 지원금액: 직급/근속기간, 지원금액, 조건 등\n- 교육과정: 과정명, 대상, 기간, 비용 등\n- 예산계획: 항목, 예산액, 비율, 비고 등\n\n사용자 키워드를 분석하여 위와 같은 구조화된 정보가 있으면 반드시 해당 내용으로 표를 만드세요. 예시:\n키워드에 '리더 5만원, 파트장 10만원, 팀장 20만원'이 있다면:\n[{\"직급/조건\": \"근속 3년이상 리더\", \"지원금액\": \"50,000원\", \"비고\": \"월 지급\"}, {\"직급/조건\": \"근속 3년이상 파트장\", \"지원금액\": \"100,000원\", \"비고\": \"월 지급\"}]\n\n표로 만들 적절한 데이터가 없다면 items 필드는 빈 배열 []로 설정하고, 모든 내용을 body에 텍스트로 작성하세요. 절대로 사용자 키워드와 무관한 예시 데이터를 사용하지 마세요. 응답은 `title`, `purpose`, `body`, `items`, `remarks` JSON 형식입니다.",
        "공지문": "당신은 한국 기업의 사내 커뮤니케이션 담당자입니다. 키워드와 첨부파일 내용을 바탕으로 '사내 공지문' 초안을 생성합니다. 'details' 필드에는 `1.`, `  1)`, `    (1)` 의 위계질서를 준수하는 번호 매기기를 사용하고, 각 문장의 마침표 후에는 반드시 줄바꿈을 해주세요. \n\n**표 생성 규칙**: 사용자 키워드에 표로 정리하면 효과적인 정보가 있다면 'items' 필드를 추가하세요. 예를 들어:\n- 일정표: 날짜, 시간, 내용, 장소\n- 교육과정: 과정명, 대상, 기간, 신청방법\n- 혜택/제도: 대상, 지원내용, 조건, 신청기한\n- 변경사항: 기존, 변경후, 시행일, 비고\n\n표가 필요한 데이터가 있으면 items 필드에 배열로 포함하고, 없으면 생략하세요. 응답은 'title', 'target', 'summary', 'details', 'contact' key와 필요시 'items' key를 포함하는 JSON 형식이어야 합니다.",
        "공문": "당신은 대외 문서를 담당하는 총무팀 직원입니다. 키워드와 첨부파일 내용을 바탕으로 격식에 맞는 '공문' 초안을 생성합니다. 본문 작성 시 `1.`, `  1)`, `    (1)` 의 위계질서를 준수하고, 각 문장의 마침표 후에는 줄바꿈을 해주세요. \n\n**표 생성 규칙**: 사용자 키워드에 표로 정리하면 효과적인 정보가 있다면 'items' 필드를 추가하세요. 예를 들어:\n- 행사일정: 일시, 장소, 내용, 참석대상\n- 제출서류: 서류명, 제출기한, 제출처, 비고\n- 협력요청: 항목, 요청사항, 기한, 담당부서\n- 비용내역: 항목, 금액, 용도, 비고\n\n표가 필요한 데이터가 있으면 items 필드에 배열로 포함하고, 없으면 생략하세요. 응답은 'sender_org', 'receiver', 'cc', 'title', 'body', 'sender_name' key와 필요시 'items' key를 포함하는 JSON 형식이어야 합니다.",
        "비즈니스 이메일": "당신은 비즈니스 커뮤니케이션 전문가입니다. 키워드와 첨부파일 내용을 바탕으로 전문적인 '비즈니스 이메일' 초안을 생성합니다. 본문 작성 시 `1.`, `  1)`, `    (1)` 의 위계질서를 준수하고, 각 문장의 마침표 후에는 줄바꿈을 해주세요. \n\n**표 생성 규칙**: 사용자 키워드에 표로 정리하면 효과적인 정보가 있다면 'items' 필드를 추가하세요. 예를 들어:\n- 미팅일정: 날짜, 시간, 안건, 참석자\n- 견적서: 항목, 수량, 단가, 금액\n- 업무일정: 업무명, 담당자, 기한, 상태\n- 제품정보: 제품명, 사양, 가격, 배송일\n\n표가 필요한 데이터가 있으면 items 필드에 배열로 포함하고, 없으면 생략하세요. 응답은 `subject`, `body`, `closing` key와 필요시 'items' key를 포함하는 JSON 형식이어야 합니다. `closing`에는 회사명, 연락처, 이메일 주소 등의 서명 정보를 포함하지 마세요. 단순히 인사말이나 마무리 문구만 포함하세요."
    }
    
    # 학습된 내용으로 프롬프트 강화
    enhanced_system_prompt = get_learning_enhanced_prompt(base_prompts[doc_type], doc_type)
    
    prompts = {
        "품의서": {"system": enhanced_system_prompt, "user": user_prompt},
        "공지문": {"system": enhanced_system_prompt, "user": user_prompt},
        "공문": {"system": enhanced_system_prompt, "user": user_prompt},
        "비즈니스 이메일": {"system": enhanced_system_prompt, "user": user_prompt}
    }
    return get_ai_response(prompts[doc_type]["system"], prompts[doc_type]["user"])

# --- 파일 읽기 및 텍스트 처리 함수들 ---
def read_uploaded_file(uploaded_file):
    if not uploaded_file:
        return ""
        
    # 파일 크기 제한 (10MB)
    max_file_size = 10 * 1024 * 1024  # 10MB
    if hasattr(uploaded_file, 'size') and uploaded_file.size > max_file_size:
        st.error(f"파일 크기가 너무 큽니다. 10MB 이하의 파일을 업로드해주세요.")
        return ""
    
    try:
        file_extension = uploaded_file.name.split('.')[-1].lower()
        
        if file_extension == "pdf":
            try:
                pdf_reader = PyPDF2.PdfReader(uploaded_file)
                if len(pdf_reader.pages) > 50:
                    st.warning("PDF 파일이 너무 깁니다. 처음 50페이지만 처리합니다.")
                
                text = ""
                for i, page in enumerate(pdf_reader.pages[:50]):
                    page_text = page.extract_text() or ""
                    text += page_text
                    
                if not text.strip():
                    st.warning("PDF에서 텍스트를 추출할 수 없습니다.")
                return text
            except Exception as e:
                st.error(f"PDF 파일 처리 중 오류: {str(e)}")
                return ""
                
        elif file_extension == "docx":
            try:
                doc = Document(uploaded_file)
                text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
                if not text.strip():
                    st.warning("Word 문서에서 텍스트를 찾을 수 없습니다.")
                return text
            except Exception as e:
                st.error(f"Word 파일 처리 중 오류: {str(e)}")
                return ""
                
        elif file_extension == "pptx":
            try:
                prs = Presentation(uploaded_file)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip(): 
                            text += shape.text + "\n"
                if not text.strip():
                    st.warning("PowerPoint에서 텍스트를 찾을 수 없습니다.")
                return text
            except Exception as e:
                st.error(f"PowerPoint 파일 처리 중 오류: {str(e)}")
                return ""
                
        elif file_extension in ['xlsx', 'xls']:
            try:
                df = pd.read_excel(uploaded_file, engine='openpyxl')
                if df.empty:
                    st.warning("Excel 파일이 비어있습니다.")
                    return ""
                return df.head(100).to_string()  # 첫 100행만 처리
            except Exception as e:
                st.error(f"Excel 파일 처리 중 오류: {str(e)}")
                return ""
                
        elif file_extension == "txt":
            try:
                content = uploaded_file.getvalue()
                text = content.decode("utf-8")
                if not text.strip():
                    st.warning("텍스트 파일이 비어있습니다.")
                return text
            except UnicodeDecodeError:
                try:
                    text = content.decode("euc-kr")
                    return text
                except UnicodeDecodeError:
                    st.error("텍스트 파일의 인코딩을 인식할 수 없습니다.")
                    return ""
            except Exception as e:
                st.error(f"텍스트 파일 처리 중 오류: {str(e)}")
                return ""
        else:
            st.warning(f"지원하지 않는 파일 형식입니다: .{file_extension}")
            return ""
            
    except Exception as e:
        st.error(f"'{uploaded_file.name}' 파일을 읽는 중 예상치 못한 오류가 발생했습니다: {str(e)}")
        return ""

def renumber_text(text):
    lines = text.split('\n')
    new_lines = []; counters = [0, 0, 0]
    for line in lines:
        stripped_line = line.lstrip()
        indentation = len(line) - len(stripped_line)
        match = re.match(r'^(\d+\.|\d+\)|\(\d+\)|\-|\*)\s+', stripped_line)
        if match:
            level = indentation // 2
            if level > 2: level = 2
            for i in range(level + 1, len(counters)): counters[i] = 0
            counters[level] += 1
            if level == 0: new_prefix = f"{counters[level]}. "
            elif level == 1: new_prefix = f"{'  ' * level}{counters[level]}) "
            else: new_prefix = f"{'  ' * level}({counters[level]}) "
            content_part = stripped_line[len(match.group(1)):].lstrip()
            new_lines.append("  " * level + new_prefix + content_part)
        else:
            new_lines.append(line)
    return "\n".join(new_lines)

def clean_text(text):
    if not isinstance(text, str): return ""
    # 마크다운 헤더 제거
    processed_text = re.sub(r'^\s*#+\s*', '', text, flags=re.MULTILINE)
    
    # 마침표 뒤에 줄바꿈 추가 (번호 매기기 제외)
    # 제외 조건:
    # 1. 숫자.공백 패턴 (1. , 2. , 3. 등)
    # 2. 공백+숫자) 패턴 (  1), (1) 등) 
    # 3. 이미 줄바꿈이 있는 경우
    # 4. 문자열 끝인 경우
    # 문장 마침표만 감지하도록 개선된 패턴
    processed_text = re.sub(r'(?<!\d)\.(?!\s*\n)(?!\s*$)(?!\s+[0-9])(?!\s*\))(?=\s*[가-힣A-Za-z])', '.\n', processed_text)
    
    # 번호 매기기 정리
    processed_text = renumber_text(processed_text)
    return processed_text

def text_to_html(text, for_email=False): 
    """텍스트를 HTML 형식으로 변환"""
    if isinstance(text, dict):
        # JSON 객체 형태로 된 경우 텍스트로 변환
        formatted_text = ""
        for key, value in text.items():
            if key.strip() in ['1.', '2.', '3.', '4.', '5.']:
                formatted_text += f"{key} {value}\n"
            elif key.strip().endswith(')') and key.strip().replace(')', '').strip().isdigit():
                formatted_text += f"  {key} {value}\n"
            elif key.strip().startswith('(') and key.strip().endswith(')'):
                formatted_text += f"    {key} {value}\n"
            else:
                formatted_text += f"{key} {value}\n"
        text = formatted_text
    
    if for_email:
        # 이메일의 경우 clean_text 처리를 하지 않고 기본 줄바꿈만 처리
        if not isinstance(text, str): 
            text = ""
        # 마크다운 헤더만 제거하고 자동 줄바꿈은 추가하지 않음
        processed_text = re.sub(r'^\s*#+\s*', '', text, flags=re.MULTILINE)
        return processed_text.replace('\n', '<br>')
    else:
        return clean_text(text).replace('\n', '<br>')

def validate_input_length(text, min_length=0, max_length=10000, field_name="입력"):
    """입력 텍스트 길이 유효성 검사"""
    if not text:
        return f"{field_name}을(를) 입력해주세요."
    
    text_length = len(text.strip())
    if text_length < min_length:
        return f"{field_name}이(가) 너무 짧습니다. 최소 {min_length}자 이상 입력해주세요."
    elif text_length > max_length:
        return f"{field_name}이(가) 너무 깁니다. {max_length}자 이하로 입력해주세요."
    
    return None

def show_progress_with_status(steps, delay=0.5):
    """진행률과 상태 메시지를 표시하는 함수"""
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    import time
    for i, step_message in enumerate(steps):
        progress = (i + 1) / len(steps)
        progress_bar.progress(progress)
        status_text.text(step_message)
        time.sleep(delay)
    
    return progress_bar, status_text

def validate_document_fields(doc_type, data):
    """문서 유형별 필드 유효성 검사"""
    errors = []
    
    if doc_type == '품의서':
        if not data.get("title") or len(data["title"].strip()) < 5:
            errors.append("제목을 5자 이상 입력해주세요.")
        if not data.get("purpose") or len(data["purpose"].strip()) < 20:
            errors.append("목적을 20자 이상 입력해주세요.")
    elif doc_type == '공지문':
        if not data.get("title") or len(data["title"].strip()) < 5:
            errors.append("제목을 5자 이상 입력해주세요.")
        if not data.get("target") or len(data["target"].strip()) < 2:
            errors.append("대상을 2자 이상 입력해주세요.")
    elif doc_type == '공문':
        if not data.get("sender_org") or len(data["sender_org"].strip()) < 3:
            errors.append("발신 기관명을 3자 이상 입력해주세요.")
        if not data.get("receiver") or len(data["receiver"].strip()) < 3:
            errors.append("수신을 3자 이상 입력해주세요.")
    elif doc_type == '비즈니스 이메일':
        if not data.get("subject") or len(data["subject"].strip()) < 5:
            errors.append("제목을 5자 이상 입력해주세요.")
        if not data.get("body") or len(data["body"].strip()) < 10:
            errors.append("본문을 10자 이상 입력해주세요.")
    
    return errors

def generate_pdf(html_content):
    font_css = CSS(string="@import url('https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@400;700&display=swap'); body { font-family: 'Noto Sans KR', sans-serif; }")
    return HTML(string=html_content).write_pdf(stylesheets=[font_css])

def generate_docx(draft_data, doc_type, signature_data={}):
    doc = Document()
    style = doc.styles['Normal']; style.font.name = '맑은 고딕'; style.font.size = Pt(11)
    if doc_type == '품의서':
        h = doc.add_heading(draft_data.get('title', '제목 없음'), level=1); h.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(clean_text(draft_data.get('purpose', '')))
        doc.add_paragraph("- 아 래 -").alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("1. 상세 내역", level=2)
        
        # 텍스트 내용 먼저 추가
        if "body" in draft_data and draft_data.get("body"):
            doc.add_paragraph(clean_text(draft_data.get('body', '')))
            if "items" in draft_data and draft_data["items"]:
                doc.add_paragraph("")  # 빈 줄 추가
        
        # 표 데이터 추가
        if "items" in draft_data and draft_data["items"]:
            df = pd.DataFrame(draft_data["items"])
            if not df.empty:
                table = doc.add_table(rows=1, cols=len(df.columns), style='Table Grid')
                hdr_cells = table.rows[0].cells
                for i, col_name in enumerate(df.columns): 
                    hdr_cells[i].text = col_name
                for _, row in df.iterrows():
                    row_cells = table.add_row().cells
                    for i, col_name in enumerate(df.columns): 
                        row_cells[i].text = str(row[col_name])
        doc.add_heading("2. 비고", level=2)
        doc.add_paragraph(clean_text(draft_data.get('remarks', '')))
        p_end = doc.add_paragraph("끝."); p_end.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    elif doc_type == '공지문':
        h = doc.add_heading(draft_data.get('title', '제목 없음'), level=1); h.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(f"대상: {draft_data.get('target', '')}")
        doc.add_paragraph(f"핵심 요약: {draft_data.get('summary', '')}")
        doc.add_paragraph("-" * 30)
        doc.add_paragraph(clean_text(draft_data.get('details', '')))
        
        # 표 데이터 추가
        if "items" in draft_data and draft_data["items"]:
            try:
                df = pd.DataFrame(draft_data["items"])
                if not df.empty:
                    doc.add_paragraph("")  # 빈 줄 추가
                    table = doc.add_table(rows=1, cols=len(df.columns), style='Table Grid')
                    hdr_cells = table.rows[0].cells
                    for i, col_name in enumerate(df.columns): 
                        hdr_cells[i].text = col_name
                    for _, row in df.iterrows():
                        row_cells = table.add_row().cells
                        for i, col_name in enumerate(df.columns): 
                            row_cells[i].text = str(row[col_name])
            except Exception as e:
                doc.add_paragraph(f"표 생성 중 오류: {str(e)}")
                
        doc.add_paragraph(f"\n문의: {draft_data.get('contact', '')}")
    elif doc_type == '공문':
        h = doc.add_heading("공 식 문 서", level=1); h.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(f"발신: {draft_data.get('sender_org', '')}")
        doc.add_paragraph(f"수신: {draft_data.get('receiver', '')}")
        doc.add_paragraph(f"참조: {draft_data.get('cc', '')}")
        doc.add_paragraph("-" * 30)
        doc.add_paragraph(f"제목: {draft_data.get('title', '')}")
        doc.add_paragraph(clean_text(draft_data.get('body', '')))
        
        # 표 데이터 추가
        if "items" in draft_data and draft_data["items"]:
            try:
                df = pd.DataFrame(draft_data["items"])
                if not df.empty:
                    doc.add_paragraph("")  # 빈 줄 추가
                    table = doc.add_table(rows=1, cols=len(df.columns), style='Table Grid')
                    hdr_cells = table.rows[0].cells
                    for i, col_name in enumerate(df.columns): 
                        hdr_cells[i].text = col_name
                    for _, row in df.iterrows():
                        row_cells = table.add_row().cells
                        for i, col_name in enumerate(df.columns): 
                            row_cells[i].text = str(row[col_name])
            except Exception as e:
                doc.add_paragraph(f"표 생성 중 오류: {str(e)}")
        
        p = doc.add_paragraph(f"\n\n{draft_data.get('sender_name', '')}"); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    elif doc_type == '비즈니스 이메일':
        doc.add_paragraph(f"받는 사람: {signature_data.get('recipient_name', '')} {signature_data.get('recipient_title', '')}")
        doc.add_paragraph(f"참조: {draft_data.get('cc', '')}")
        doc.add_paragraph(f"제목: {draft_data.get('subject', '')}")
        doc.add_paragraph("-" * 30)
        doc.add_paragraph(f"안녕하세요, {signature_data.get('recipient_name', '')} {signature_data.get('recipient_title', '')}님.")
        doc.add_paragraph(f"{signature_data.get('signature_name', '')} {signature_data.get('signature_title', '')}입니다.")
        doc.add_paragraph() 
        doc.add_paragraph(clean_text(draft_data.get('body', '')))
        
        # 표 데이터 추가
        if "items" in draft_data and draft_data["items"]:
            try:
                df = pd.DataFrame(draft_data["items"])
                if not df.empty:
                    doc.add_paragraph("")  # 빈 줄 추가
                    table = doc.add_table(rows=1, cols=len(df.columns), style='Table Grid')
                    hdr_cells = table.rows[0].cells
                    for i, col_name in enumerate(df.columns): 
                        hdr_cells[i].text = col_name
                    for _, row in df.iterrows():
                        row_cells = table.add_row().cells
                        for i, col_name in enumerate(df.columns): 
                            row_cells[i].text = str(row[col_name])
            except Exception as e:
                doc.add_paragraph(f"표 생성 중 오류: {str(e)}")
        
        doc.add_paragraph(clean_text(draft_data.get('closing', '')))
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

st.set_page_config(page_title="문서 작성 도우미", layout="wide")
env = Environment(loader=FileSystemLoader('.'))
def load_template(template_name): return env.get_template(template_name)
def generate_html(template, context): return template.render(context)

def clear_all_state():
    """문서 유형 변경 시 관련 상태만 초기화"""
    keys_to_keep = ['doc_type_selector']
    keys_to_remove = [key for key in st.session_state.keys() if key not in keys_to_keep]
    for key in keys_to_remove:
        del st.session_state[key]

st.sidebar.title("📑 문서 종류 선택")
# 이전 문서 타입 저장
if 'previous_doc_type' not in st.session_state:
    st.session_state.previous_doc_type = None

doc_type = st.sidebar.radio("작성할 문서의 종류를 선택하세요.", ('품의서', '공지문', '공문', '비즈니스 이메일'), key="doc_type_selector")

# --- 설정 섹션 ---
st.sidebar.divider()
st.sidebar.title("⚙️ 설정")

# AI 모델 선택
st.sidebar.subheader("🤖 AI 모델 설정")
current_model = st.session_state.selected_model
st.sidebar.info(f"현재 모델: **{current_model}**")

# 모델 비용 정보 표시
model_costs = {
    "gpt-4o-mini": "💚 저렴 (기본)",
    "gpt-4o": "💰 비쌈 (고성능)",
    "gpt-4-turbo": "💸 매우 비쌈", 
    "gpt-3.5-turbo": "💚 매우 저렴"
}
st.sidebar.caption(f"비용: {model_costs.get(current_model, '알 수 없음')}")

# 모델 변경 요청 처리
if st.sidebar.button("🔧 모델 변경하기", use_container_width=True):
    if not st.session_state.model_password_verified:
        # 비밀번호 입력 상태로 변경
        if 'show_password_input' not in st.session_state:
            st.session_state.show_password_input = True
        else:
            st.session_state.show_password_input = not st.session_state.show_password_input

# 비밀번호 입력 화면
if st.session_state.get('show_password_input', False) and not st.session_state.model_password_verified:
    password = st.sidebar.text_input("🔐 비밀번호 입력", type="password", placeholder="모델 변경 비밀번호")
    
    col1, col2 = st.sidebar.columns(2)
    with col1:
        if st.button("확인", use_container_width=True):
            if password == "admin123":  # 비밀번호를 여기서 설정 (변경 가능)
                st.session_state.model_password_verified = True
                st.session_state.show_password_input = False
                st.sidebar.success("✅ 인증 성공!")
                st.rerun()
            else:
                st.sidebar.error("❌ 잘못된 비밀번호입니다.")
    
    with col2:
        if st.button("취소", use_container_width=True):
            st.session_state.show_password_input = False
            st.rerun()

# 인증된 경우 모델 선택 표시
if st.session_state.model_password_verified:
    st.sidebar.subheader("모델 선택")
    new_model = st.sidebar.selectbox(
        "사용할 모델을 선택하세요:",
        ["gpt-4o-mini", "gpt-4o", "gpt-4-turbo", "gpt-3.5-turbo"],
        index=["gpt-4o-mini", "gpt-4o", "gpt-4-turbo", "gpt-3.5-turbo"].index(current_model)
    )
    
    if st.sidebar.button("💾 모델 저장", use_container_width=True):
        st.session_state.selected_model = new_model
        st.session_state.model_password_verified = False
        st.sidebar.success(f"✅ 모델이 **{new_model}**로 변경되었습니다!")
        st.rerun()
    
    if st.sidebar.button("❌ 취소", use_container_width=True):
        st.session_state.model_password_verified = False
        st.rerun()

st.sidebar.divider()

# 학습 상태 표시 (간단하게)
if learning_status["manual"] or learning_status["samples"] or learned_documents.get('files'):
    if learned_documents.get('files'):
        # 새로운 files 구조가 있는 경우
        files_data = learned_documents.get('files', {})
        successful_files = [f for f, data in files_data.items() if data.get('success')]
        total_files = len(files_data)
        
        st.sidebar.success("📚 PDF 학습 완료!")
        st.sidebar.caption(f"총 {total_files}개 파일 중 {len(successful_files)}개 성공")
        
        summary = learned_documents.get('summary', {})
        if summary:
            total_length = summary.get('total_content_length', 0)
            st.sidebar.caption(f"학습된 내용: {total_length:,}자")
    else:
        # 기존 방식
        st.sidebar.success("📚 학습 완료!")
        summary = learned_documents.get('summary', {})
        if summary:
            total_length = summary.get('total_content_length', 0)
            st.sidebar.caption(f"학습된 내용: {total_length:,}자")
    
    learned_at = learned_documents.get('learned_at', '알 수 없음')
    st.sidebar.caption(f"학습 일시: {learned_at}")
else:
    st.sidebar.warning("📖 아직 학습되지 않음")

# 학습 실행 버튼
if st.sidebar.button("📚 PDF 문서 학습하기", use_container_width=True):
    try:
        with st.spinner("PDF 문서를 학습 중입니다..."):
            # 실제 PDF 파일 읽기
            from datetime import datetime
            
            def read_pdf_file(filename):
                """PDF 파일을 읽어서 텍스트를 추출합니다."""
                try:
                    if not os.path.exists(filename):
                        return f"파일 '{filename}'을 찾을 수 없습니다.", False
                    
                    # PyPDF2 import 확인
                    try:
                        import PyPDF2
                    except ImportError:
                        return f"PyPDF2 모듈을 찾을 수 없습니다. PDF 읽기 기능이 비활성화됩니다.", False
                    
                    with open(filename, 'rb') as file:
                        pdf_reader = PyPDF2.PdfReader(file)
                        text = ""
                        page_count = len(pdf_reader.pages)
                        
                        for i, page in enumerate(pdf_reader.pages):
                            try:
                                page_text = page.extract_text()
                                if page_text:
                                    text += page_text + "\n"
                            except Exception as page_error:
                                st.warning(f"⚠️ {filename} 페이지 {i+1} 읽기 실패: {str(page_error)}")
                        
                        if not text.strip():
                            return f"PDF '{filename}'에서 텍스트를 추출할 수 없습니다. (총 {page_count}페이지)", False
                        
                        return text.strip(), True
                        
                except Exception as e:
                    return f"PDF '{filename}' 읽기 중 오류: {str(e)}", False
            
            # 폴더에서 모든 PDF 파일 자동 검색
            import glob
            pdf_files = glob.glob('*.pdf') + glob.glob('*.PDF')
            
            st.info(f"폴더에서 {len(pdf_files)}개의 PDF 파일을 발견했습니다.")
            
            learned_content = {
                'learned_at': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'status': 'learned',
                'files': {},
                'summary': {
                    'total_files': len(pdf_files),
                    'successful_files': 0,
                    'failed_files': 0,
                    'total_content_length': 0
                }
            }
            
            # 각 PDF 파일을 순차적으로 학습
            for i, pdf_file in enumerate(pdf_files):
                st.info(f"📖 {pdf_file} 읽는 중... ({i+1}/{len(pdf_files)})")
                
                try:
                    content, success = read_pdf_file(pdf_file)
                    
                    if success:
                        learned_content['files'][pdf_file] = {
                            'filename': pdf_file,
                            'content': content,
                            'source': 'pdf_extracted',
                            'length': len(content),
                            'success': True
                        }
                        learned_content['summary']['successful_files'] += 1
                        learned_content['summary']['total_content_length'] += len(content)
                        st.success(f"✅ {pdf_file} 학습 완료 ({len(content):,}자)")
                    else:
                        learned_content['files'][pdf_file] = {
                            'filename': pdf_file,
                            'content': content,
                            'source': 'error',
                            'length': 0,
                            'success': False
                        }
                        learned_content['summary']['failed_files'] += 1
                        st.error(f"❌ {pdf_file}: {content[:100]}...")
                        
                except Exception as e:
                    learned_content['files'][pdf_file] = {
                        'filename': pdf_file,
                        'content': f"처리 중 오류 발생: {str(e)}",
                        'source': 'error',
                        'length': 0,
                        'success': False
                    }
                    learned_content['summary']['failed_files'] += 1
                    st.error(f"❌ {pdf_file} 처리 실패: {str(e)}")
            
            # 기존 파일들 호환성 유지 (manual, samples 키 생성)
            manual_files = [f for f in pdf_files if '메뉴얼' in f or 'manual' in f.lower()]
            samples_files = [f for f in pdf_files if '품의서' in f or '모음' in f or 'sample' in f.lower()]
            
            if manual_files:
                learned_content['manual'] = learned_content['files'][manual_files[0]]
            else:
                learned_content['manual'] = {
                    'content': "기본 가이드라인을 사용합니다.",
                    'source': 'fallback_guidelines',
                    'success': False
                }
            
            if samples_files:
                learned_content['samples'] = learned_content['files'][samples_files[0]]
            else:
                learned_content['samples'] = {
                    'content': "기본 샘플 패턴을 사용합니다.",
                    'source': 'fallback_patterns', 
                    'success': False
                }
            
            # 학습 결과 확인 및 저장
            successful_files = learned_content['summary']['successful_files']
            total_files = learned_content['summary']['total_files']
            
            if successful_files > 0:
                # 성공한 파일이 있는 경우에만 저장
                with open('learned_documents.json', 'w', encoding='utf-8') as f:
                    json.dump(learned_content, f, ensure_ascii=False, indent=2)
                
                st.success(f"📚 PDF 학습 완료! 총 {total_files}개 파일 중 {successful_files}개 성공")
                st.info(f"학습된 내용: {learned_content['summary']['total_content_length']:,}자")
                st.info(f"학습 일시: {learned_content['learned_at']}")
                
                # 학습 완료 후 다시 로드
                if load_learned_documents():
                    st.sidebar.success("✅ PDF 학습이 완료되었습니다!")
                    st.rerun()
                else:
                    st.sidebar.error("❌ 학습 결과를 로드할 수 없습니다.")
            else:
                # 성공한 파일이 없는 경우
                st.error(f"❌ PDF 학습 실패! 총 {total_files}개 파일 모두 읽기 실패")
                st.warning("PyPDF2 모듈이나 PDF 파일에 문제가 있을 수 있습니다.")
                
                # 실패 상세 정보 표시
                for pdf_file, file_data in learned_content['files'].items():
                    if not file_data['success']:
                        st.error(f"📄 {pdf_file}: {file_data['content'][:200]}...")
                
                st.info("💡 해결 방법: requirements.txt에 PyPDF2가 포함되어 있는지 확인하고, Streamlit을 재시작해보세요.")
                
    except Exception as e:
        st.sidebar.error(f"❌ 학습 실행 중 오류: {str(e)}")

# 학습 상태 초기화 버튼
if learning_status["manual"] or learning_status["samples"]:
    if st.sidebar.button("🗑️ 학습 데이터 초기화", use_container_width=True):
        if reset_learning_data():
            st.sidebar.success("✅ 학습 데이터가 초기화되었습니다!")
            st.rerun()

# 문서 타입이 변경된 경우에만 상태 초기화
if st.session_state.previous_doc_type != doc_type:
    clear_all_state()
    st.session_state.previous_doc_type = doc_type

# 세션 상태 초기화 - 키 생성 방식 개선
draft_key = f"draft_{doc_type.replace(' ', '_')}"
html_key = f"html_{doc_type.replace(' ', '_')}"

# 필요한 상태만 초기화
state_defaults = {
    draft_key: {},
    html_key: "",
    "clarifying_questions": None,
    "current_keywords": "",
    "file_processing_complete": False,
    "ai_generation_complete": False
}

for key, default_value in state_defaults.items():
    if key not in st.session_state:
        st.session_state[key] = default_value

if openai_available:
    st.title(f"✍️ {doc_type} 작성 가이드")
    col1, col2 = st.columns([3, 1])
    with col1:
        st.success("🤖 AI 기능이 활성화되었습니다!")
    with col2:
        if learning_status["manual"] or learning_status["samples"]:
            st.success("📚 학습 완료")
        else:
            st.info("📖 미학습")
else:
    st.title(f"📝 {doc_type} 템플릿")
    st.error("⚠️ AI 기능이 비활성화되었습니다. OpenAI API 키를 설정해주세요.")

if not st.session_state.clarifying_questions:
    if openai_available:
        if not (learning_status["manual"] or learning_status["samples"] or learned_documents.get('files')):
            st.info("💡 **팁**: 사이드바에서 'PDF 문서 학습하기'를 클릭하면 더욱 전문적인 문서를 생성할 수 있습니다.")
    else:
        st.markdown("현재 AI 기능이 비활성화되어 있습니다. OpenAI API 키를 설정하면 자동 문서 생성 기능을 사용할 수 있습니다.")
        with st.expander("API 키 설정 방법"):
            st.markdown("""
            1. [OpenAI 웹사이트](https://platform.openai.com/)에서 API 키를 발급받으세요
            2. Streamlit Cloud의 앱 설정에서 Secrets 섹션으로 이동하세요
            3. 다음과 같이 API 키를 추가하세요:
            ```
            OPENAI_API_KEY = "your-api-key-here"
            ```
            4. 앱을 재시작하세요
            """)
    sub_type = ""
    if doc_type == "품의서":
        sub_type = st.selectbox("품의서 세부 유형을 선택하세요:", ["선택 안함", "비용 집행", "신규 사업/계약", "인사/정책 변경", "결과/사건 보고"])
    keywords = st.text_area("핵심 키워드", placeholder="예: 영업팀 태블릿 5대 구매, 총 예산 400만원, 업무용", height=100, key="keyword_input")
    
    # 입력 검증 및 안내
    if keywords:
        word_count = len(keywords.split())
        char_count = len(keywords)
        
        if char_count < 10:
            st.warning("⚠️ 너무 짧습니다. 더 상세한 내용을 입력해주세요. (최소 10자 이상)")
        elif char_count > 1000:
            st.warning("⚠️ 너무 깁니다. 1000자 이하로 입력해주세요.")
        else:
            st.success(f"✅ 적절한 길이입니다. (단어: {word_count}개, 문자: {char_count}자)")
    uploaded_files = st.file_uploader("참고 파일 업로드 (선택 사항)", type=['pdf', 'docx', 'pptx', 'xlsx', 'xls', 'txt'], accept_multiple_files=True)
    
    # 파일 업로드 안내
    if uploaded_files:
        if len(uploaded_files) > 5:
            st.error("⚠️ 최대 5개의 파일만 업로드 할 수 있습니다.")
            uploaded_files = uploaded_files[:5]
        
        total_size = sum(getattr(f, 'size', 0) for f in uploaded_files)
        if total_size > 50 * 1024 * 1024:  # 50MB 제한
            st.error("⚠️ 전체 파일 크기가 50MB를 초과합니다.")
        else:
            st.info(f"파일 {len(uploaded_files)}개 업로드됨 (전체 크기: {total_size/1024/1024:.1f}MB)")
    use_clarifying_questions = st.checkbox("AI에게 추가 질문을 받아 문서 완성도 높이기 (선택 사항)")
    ai_button_disabled = not openai_available
    if ai_button_disabled:
        st.warning("⚠️ OpenAI API 키가 필요합니다. Streamlit Secrets에 OPENAI_API_KEY를 설정해주세요.")
    
    if st.button("AI 초안 생성 시작", type="primary", use_container_width=True, disabled=ai_button_disabled):
        # 입력 유효성 검사
        validation_errors = []
        
        if not keywords or len(keywords.strip()) < 10:
            validation_errors.append("핵심 키워드를 10자 이상 입력해주세요.")
        
        if len(keywords) > 1000:
            validation_errors.append("키워드는 1000자 이하로 입력해주세요.")
        
        if uploaded_files and len(uploaded_files) > 5:
            validation_errors.append("참고 파일은 최대 5개까지만 업로드 가능합니다.")
        
        if validation_errors:
            for error in validation_errors:
                st.error(f"⚠️ {error}")
        else:
            full_keywords = f"유형: {sub_type} / 내용: {keywords}" if sub_type != "선택 안함" else keywords
            st.session_state.current_keywords = full_keywords
            file_context = ""
            
            # 파일 처리 진행률 표시
            if uploaded_files:
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                for i, uploaded_file in enumerate(uploaded_files):
                    progress = (i + 1) / len(uploaded_files)
                    progress_bar.progress(progress)
                    status_text.text(f"파일 처리 중: {uploaded_file.name} ({i+1}/{len(uploaded_files)})")
                    
                    file_text = read_uploaded_file(uploaded_file)
                    if file_text:
                        file_context += f"--- 첨부 파일: {uploaded_file.name} ---\n{file_text}\n\n"
                
                progress_bar.empty()
                status_text.empty()
                st.success(f"파일 처리 완료: {len(uploaded_files)}개 파일")
            
            analysis_complete = True
            if use_clarifying_questions:
                with st.spinner("🤖 AI가 키워드를 분석하여 추가 질문을 준비 중입니다..."):
                    analysis = analyze_keywords(full_keywords, doc_type)
                    if analysis and analysis.get("status") == "incomplete":
                        st.session_state.clarifying_questions = analysis.get("questions", [])
                        analysis_complete = False
                        st.info("🔍 문서 품질 향상을 위해 추가 정보가 필요합니다.")
                        st.rerun()
            if analysis_complete:
                # AI 생성 진행률 표시
                steps = [
                    "🤖 AI가 문서 구조를 분석하고 있습니다...",
                    f"📝 {doc_type} 컨텐츠를 생성하고 있습니다...",
                    "✨ 최종 검토 및 포맷팅 중입니다..."
                ]
                progress_bar, status_text = show_progress_with_status(steps)
                
                ai_result = generate_ai_draft(doc_type, full_keywords, file_context)
                
                progress_bar.progress(1.0)
                status_text.text("✅ 문서 생성 완료!")
                import time
                time.sleep(1)
                
                progress_bar.empty()
                status_text.empty()
                    
                if ai_result:
                    st.session_state[draft_key] = ai_result
                    st.session_state[html_key] = ""
                    st.success("✨ AI가 문서 초안을 성공적으로 생성했습니다! 아래에서 내용을 확인하고 수정해주세요.")
                else:
                    st.error("문서 생성에 실패했습니다. 다시 시도해주세요.")
        
    # 추가 도움말 제공
    with st.expander("효과적인 키워드 작성 팁"):
        st.markdown("""
        **좋은 키워드 예시:**
        - "마케팅팀 노트북 10대 구매, 예산 500만원, 2024년 4분기 지급"
        - "신입사원 원격근무 제도 도입, 2025년 1월부터 시행"
        - "고객서비스 운영시간 연장, 평일 21시까지, 인력 증원 필요"
        
        **피해야 할 키워드:**
        - 너무 간단: "노트북 구매"
        - 너무 모호: "여러 가지 사무용품 구매 관련"
        - 배경 설명 없이: "예산 승인 요청"
        """)
else:
    st.subheader("AI의 추가 질문 🙋‍♂️")
    st.info("문서의 완성도를 높이기 위해 몇 가지 추가 정보가 필요합니다.")
    answers = {}
    for i, q in enumerate(st.session_state.clarifying_questions):
        answer = st.text_input(q, key=f"q_{i}")
        answers[q] = answer
        
        # 질문별 입력 검증
        if answer and len(answer.strip()) < 3:
            st.warning(f"⚠️ 질문 {i+1}: 너무 짧습니다. 더 상세히 답변해주세요.")
        elif answer and len(answer) > 500:
            st.warning(f"⚠️ 질문 {i+1}: 너무 깁니다. 500자 이하로 입력해주세요.")
    if st.button("답변 제출하고 문서 생성하기", type="primary", use_container_width=True, disabled=not openai_available):
        # 답변 유효성 검사
        answered_questions = [q for q, a in answers.items() if a.strip()]
        if len(answered_questions) == 0:
            st.warning("⚠️ 적어도 하나의 질문에 답변해주세요.")
        else:
            combined_info = st.session_state.current_keywords + "\n[추가 정보]\n"
            for q, a in answers.items():
                if a: combined_info += f"- {q}: {a}\n"
            
            # 진행률 표시
            steps = [
                "🔍 추가 정보를 분석하고 있습니다...",
                f"📝 향상된 {doc_type}를 생성하고 있습니다...",
                "✨ 최종 검토 중입니다..."
            ]
            progress_bar, status_text = show_progress_with_status(steps)
            
            ai_result = generate_ai_draft(doc_type, combined_info)
            
            progress_bar.progress(1.0)
            status_text.text("✅ 개선된 문서 생성 완료!")
            import time
            time.sleep(1)
            
            progress_bar.empty()
            status_text.empty()
            
            if ai_result:
                st.session_state[draft_key] = ai_result
                st.session_state.clarifying_questions = None
                st.session_state.current_keywords = ""
                st.session_state[html_key] = ""
                st.success("✨ 추가 정보를 반영한 개선된 문서가 생성되었습니다!")
                st.rerun()
            else:
                st.error("문서 생성에 실패했습니다. 다시 시도해주세요.")

st.divider()
draft = st.session_state.get(draft_key, {})

if draft:
    preview_button = False; signature_data = {}
    st.markdown("---")
    st.subheader("📄 AI 생성 초안 검토 및 수정")
    if doc_type == '품의서':
        p_data = draft
        title_input = st.text_input("제목", value=p_data.get("title", ""), help="결재자가 제목만 보고도 내용을 파악할 수 있도록 작성합니다.")
        if title_input and len(title_input.strip()) < 5:
            st.warning("⚠️ 제목이 너무 짧습니다. 더 드립적으로 작성해주세요.")
        elif title_input and len(title_input) > 100:
            st.warning("⚠️ 제목이 너무 깁니다. 100자 이하로 작성해주세요.")
        p_data["title"] = title_input
        
        purpose_input = st.text_area("목적 및 개요", value=p_data.get("purpose", ""), height=100, help="이 품의를 올리는 이유와 목표를 명확하고 간결하게 기술합니다. (Why)")
        if purpose_input and len(purpose_input.strip()) < 20:
            st.warning("⚠️ 목적이 너무 짧습니다. 더 상세하게 설명해주세요.")
        p_data["purpose"] = purpose_input
        
        # 텍스트 내용 편집
        st.markdown("**상세 설명 (텍스트)**")
        p_data["body_edited"] = st.text_area("배경 및 설명", value=p_data.get("body", ""), height=150, help="배경, 필요성, 추진 방법 등을 텍스트로 상세히 설명합니다.")
        
        # 표 데이터 편집
        st.markdown("**상세 내역 (표)**")
        st.caption("구체적인 항목, 수량, 금액 등을 표로 정리합니다.")
        try:
            if "items" in p_data and p_data["items"] and len(p_data["items"]) > 0:
                # AI가 생성한 표가 있는 경우 - 안전하게 DataFrame 생성
                items_data = p_data.get("items", [])
                if isinstance(items_data, list) and len(items_data) > 0:
                    # 첫 번째 항목이 딕셔너리인지 확인
                    if isinstance(items_data[0], dict):
                        try:
                            p_data["df"] = pd.DataFrame(items_data)
                            p_data["df_edited"] = st.data_editor(p_data["df"], num_rows="dynamic")
                        except Exception as e:
                            st.warning(f"⚠️ AI 생성 표 데이터에 문제가 있어 기본 형식을 사용합니다: {str(e)}")
                            # 기본 구조로 대체
                            default_items = [
                                {"항목": "노트북", "수량": "10", "단가": "500,000", "금액": "5,000,000", "비고": "마케팅팀용"}
                            ]
                            p_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
                    else:
                        # 데이터 형식이 올바르지 않은 경우
                        default_items = [
                            {"항목": "노트북", "수량": "10", "단가": "500,000", "금액": "5,000,000", "비고": "마케팅팀용"}
                        ]
                        p_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
                else:
                    # 빈 데이터인 경우
                    default_items = [
                        {"항목": "노트북", "수량": "10", "단가": "500,000", "금액": "5,000,000", "비고": "마케팅팀용"}
                    ]
                    p_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
            else:
                # 표가 없는 경우 기본 구조 제공
                default_items = [
                    {"항목": "노트북", "수량": "10", "단가": "500,000", "금액": "5,000,000", "비고": "마케팅팀용"}
                ]
                p_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
        except Exception as e:
            st.error(f"⚠️ 표 데이터 처리 중 오류가 발생했습니다: {str(e)}")
            # 최종 fallback
            default_items = [
                {"항목": "예시 항목", "수량": "1", "단가": "100,000", "금액": "100,000", "비고": "설명"}
            ]
            p_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
        
        p_data["remarks"] = st.text_area("비고", value=p_data.get("remarks", ""), height=150, help="예상 비용(How much), 소요 기간(How long), 기대 효과 등 의사결정에 필요한 추가 정보를 기입합니다.")
        
        # 품의서 유효성 검사
        validation_errors = validate_document_fields(doc_type, p_data)
        
        if validation_errors:
            for error in validation_errors:
                st.error(f"⚠️ {error}")
            preview_button = st.button("미리보기 생성", use_container_width=True, disabled=True)
        else:
            preview_button = st.button("미리보기 생성", use_container_width=True)
    elif doc_type == '공지문':
        g_data = draft
        g_data["title"] = st.text_input("제목", value=g_data.get("title", ""), help="공지의 내용을 한눈에 파악할 수 있도록 작성합니다.")
        g_data["target"] = st.text_input("대상", value=g_data.get("target", ""), help="공지의 적용 범위를 명확히 합니다. (예: 전 직원)")
        g_data["summary"] = st.text_area("핵심 요약", value=g_data.get("summary", ""), height=100, help="본문 상단에 한두 문장으로 공지의 핵심을 요약합니다.")
        # 상세 내용이 JSON 객체 형태인 경우 텍스트로 변환
        details_value = g_data.get("details", "")
        if isinstance(details_value, dict):
            formatted_details = ""
            for key, value in details_value.items():
                if key.strip() in ['1.', '2.', '3.', '4.', '5.']:
                    formatted_details += f"{key} {value}\n"
                elif key.strip().endswith(')') and key.strip().replace(')', '').strip().isdigit():
                    formatted_details += f"  {key} {value}\n"
                elif key.strip().startswith('(') and key.strip().endswith(')'):
                    formatted_details += f"    {key} {value}\n"
                else:
                    formatted_details += f"{key} {value}\n"
            details_value = formatted_details
        
        g_data["details"] = st.text_area("상세 내용", value=details_value, height=200, help="5W1H 원칙에 따라 구체적인 정보를 제공합니다. 번호 매기기: 1. → 1) → (1)")
        
        # 표 데이터 편집 (공지문용)
        st.markdown("**상세 내역 (표) - 선택사항**")
        st.caption("일정, 교육과정, 제도 변경사항 등을 표로 정리할 수 있습니다.")
        try:
            if "items" in g_data and g_data["items"] and len(g_data["items"]) > 0:
                # AI가 생성한 표가 있는 경우
                items_data = g_data.get("items", [])
                if isinstance(items_data, list) and len(items_data) > 0 and isinstance(items_data[0], dict):
                    try:
                        g_data["df"] = pd.DataFrame(items_data)
                        g_data["df_edited"] = st.data_editor(g_data["df"], num_rows="dynamic")
                    except Exception as e:
                        st.warning(f"⚠️ AI 생성 표 데이터에 문제가 있어 기본 형식을 사용합니다: {str(e)}")
                        default_items = [
                            {"항목": "교육과정", "날짜": "2025-01-15", "시간": "09:00", "장소": "대회의실"}
                        ]
                        g_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
                else:
                    default_items = [
                        {"항목": "교육과정", "날짜": "2025-01-15", "시간": "09:00", "장소": "대회의실"}
                    ]
                    g_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
            else:
                # 표가 없는 경우 기본 구조 제공 (필요시만)
                if st.checkbox("표 추가하기 (일정, 교육과정 등)", key="add_table_gongji"):
                    default_items = [
                        {"항목": "교육과정", "날짜": "2025-01-15", "시간": "09:00", "장소": "대회의실"}
                    ]
                    g_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
                else:
                    g_data["df_edited"] = None
        except Exception as e:
            st.error(f"⚠️ 표 데이터 처리 중 오류가 발생했습니다: {str(e)}")
            g_data["df_edited"] = None
        
        g_data["contact"] = st.text_input("문의처", value=g_data.get("contact", ""), help="관련 질문에 답변할 담당자 정보입니다.")
        preview_button = st.button("미리보기 생성", use_container_width=True)
    elif doc_type == '공문':
        gm_data = draft
        gm_data["sender_org"] = st.text_input("발신 기관명", value=gm_data.get("sender_org", ""))
        gm_data["receiver"] = st.text_input("수신", value=gm_data.get("receiver", ""))
        gm_data["cc"] = st.text_input("참조", value=gm_data.get("cc", ""))
        gm_data["title"] = st.text_input("제목", value=gm_data.get("title", ""))
        gm_data["body"] = st.text_area("내용", value=gm_data.get("body", ""), height=250)
        
        # 표 데이터 편집 (공문용)
        st.markdown("**상세 내역 (표) - 선택사항**")
        st.caption("행사일정, 제출서류, 협력요청 등을 표로 정리할 수 있습니다.")
        try:
            if "items" in gm_data and gm_data["items"] and len(gm_data["items"]) > 0:
                # AI가 생성한 표가 있는 경우
                items_data = gm_data.get("items", [])
                if isinstance(items_data, list) and len(items_data) > 0 and isinstance(items_data[0], dict):
                    try:
                        gm_data["df"] = pd.DataFrame(items_data)
                        gm_data["df_edited"] = st.data_editor(gm_data["df"], num_rows="dynamic")
                    except Exception as e:
                        st.warning(f"⚠️ AI 생성 표 데이터에 문제가 있어 기본 형식을 사용합니다: {str(e)}")
                        default_items = [
                            {"항목": "제출서류", "서류명": "사업자등록증", "제출기한": "2025-01-31", "제출처": "총무팀"}
                        ]
                        gm_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
                else:
                    default_items = [
                        {"항목": "제출서류", "서류명": "사업자등록증", "제출기한": "2025-01-31", "제출처": "총무팀"}
                    ]
                    gm_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
            else:
                # 표가 없는 경우 기본 구조 제공 (필요시만)
                if st.checkbox("표 추가하기 (일정, 서류, 협력요청 등)", key="add_table_gongmun"):
                    default_items = [
                        {"항목": "제출서류", "서류명": "사업자등록증", "제출기한": "2025-01-31", "제출처": "총무팀"}
                    ]
                    gm_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
                else:
                    gm_data["df_edited"] = None
        except Exception as e:
            st.error(f"⚠️ 표 데이터 처리 중 오류가 발생했습니다: {str(e)}")
            gm_data["df_edited"] = None
        
        gm_data["sender_name"] = st.text_input("발신 명의", value=gm_data.get("sender_name", ""))
        preview_button = st.button("미리보기 생성", use_container_width=True)
    elif doc_type == '비즈니스 이메일':
        e_data = draft
        st.subheader("받는 사람 정보")
        signature_data["recipient_name"] = st.text_input("받는 사람 이름", value=e_data.get("recipient_name", ""))
        signature_data["recipient_title"] = st.text_input("받는 사람 직책", value=e_data.get("recipient_title", ""))
        e_data["cc"] = st.text_input("참조 (CC)", value=e_data.get("cc", ""))
        st.subheader("메일 내용")
        e_data["subject"] = st.text_input("제목", value=e_data.get("subject", ""))
        e_data["body"] = st.text_area("본론", value=e_data.get("body", ""), height=200)
        
        # 표 데이터 편집 (비즈니스 이메일용)
        st.markdown("**상세 내역 (표) - 선택사항**")
        st.caption("미팅일정, 견적서, 업무일정 등을 표로 정리할 수 있습니다.")
        try:
            if "items" in e_data and e_data["items"] and len(e_data["items"]) > 0:
                # AI가 생성한 표가 있는 경우
                items_data = e_data.get("items", [])
                if isinstance(items_data, list) and len(items_data) > 0 and isinstance(items_data[0], dict):
                    try:
                        e_data["df"] = pd.DataFrame(items_data)
                        e_data["df_edited"] = st.data_editor(e_data["df"], num_rows="dynamic")
                    except Exception as e:
                        st.warning(f"⚠️ AI 생성 표 데이터에 문제가 있어 기본 형식을 사용합니다: {str(e)}")
                        default_items = [
                            {"항목": "미팅일정", "날짜": "2025-01-15", "시간": "14:00", "안건": "프로젝트 계획 논의"}
                        ]
                        e_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
                else:
                    default_items = [
                        {"항목": "미팅일정", "날짜": "2025-01-15", "시간": "14:00", "안건": "프로젝트 계획 논의"}
                    ]
                    e_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
            else:
                # 표가 없는 경우 기본 구조 제공 (필요시만)
                if st.checkbox("표 추가하기 (일정, 견적, 업무 등)", key="add_table_email"):
                    default_items = [
                        {"항목": "미팅일정", "날짜": "2025-01-15", "시간": "14:00", "안건": "프로젝트 계획 논의"}
                    ]
                    e_data["df_edited"] = st.data_editor(pd.DataFrame(default_items), num_rows="dynamic")
                else:
                    e_data["df_edited"] = None
        except Exception as e:
            st.error(f"⚠️ 표 데이터 처리 중 오류가 발생했습니다: {str(e)}")
            e_data["df_edited"] = None
        
        e_data["closing"] = st.text_area("결론", value=e_data.get("closing", ""), height=100)
        with st.expander("내 서명 정보 입력/수정"):
            signature_data["signature_name"] = st.text_input("이름", value="홍길동")
            signature_data["signature_title"] = st.text_input("직책", value="대리")
            signature_data["signature_team"] = st.text_input("부서/팀", value="마케팅팀")
            signature_data["signature_phone"] = st.text_input("연락처", value="010-1234-5678")
        preview_button = st.button("이메일 본문 생성", use_container_width=True)
    
    if preview_button:
        if doc_type == '품의서':
            # 제목, 목적, 비고 업데이트
            draft['title'] = p_data["title"]
            draft['purpose'] = p_data["purpose"] 
            draft['remarks'] = p_data["remarks"]
            
            # 텍스트 내용 항상 포함
            draft['body'] = p_data["body_edited"]
            
            # 표 데이터 항상 포함 (비어있지 않은 경우에만)
            try:
                if "df_edited" in p_data and p_data["df_edited"] is not None and not p_data["df_edited"].empty:
                    # 빈 행 제거
                    filtered_df = p_data["df_edited"].dropna(how='all')
                    if not filtered_df.empty:
                        draft['items'] = filtered_df.to_dict('records')
                    else:
                        draft['items'] = []
                else:
                    draft['items'] = []
            except Exception as e:
                st.warning(f"⚠️ 표 데이터 처리 중 문제가 발생했습니다: {str(e)}")
                draft['items'] = []
            
            # 템플릿 컨텍스트 구성
            context = { 
                "title": draft["title"], 
                "purpose": text_to_html(draft["purpose"]), 
                "remarks": text_to_html(draft["remarks"]), 
                "generation_date": datetime.now().strftime('%Y-%m-%d') 
            }
            
            # 텍스트 내용 추가
            if draft.get("body"):
                context["body"] = text_to_html(draft["body"])
            
            # 표 데이터 추가
            if draft.get("items"):
                try:
                    if "df_edited" in p_data and p_data["df_edited"] is not None and not p_data["df_edited"].empty:
                        context["table_headers"] = list(p_data["df_edited"].columns)
                        context["items"] = draft["items"]
                    else:
                        context["items"] = []
                except Exception as e:
                    st.warning(f"⚠️ 표 헤더 처리 중 문제가 발생했습니다: {str(e)}")
                    context["items"] = []
            
            template = load_template('pumui_template_final.html')
            st.session_state[html_key] = generate_html(template, context)
        elif doc_type == '공지문':
            draft = g_data
            context = { "title": draft["title"], "target": draft["target"], "summary": text_to_html(draft["summary"]), "details": text_to_html(draft["details"]), "contact": draft["contact"], "generation_date": datetime.now().strftime('%Y. %m. %d.') }
            
            # 표 데이터 처리 (AI 생성 또는 사용자 편집)
            try:
                if "df_edited" in g_data and g_data["df_edited"] is not None and not g_data["df_edited"].empty:
                    # 사용자가 편집한 표 데이터 사용
                    filtered_df = g_data["df_edited"].dropna(how='all')
                    if not filtered_df.empty:
                        context["table_headers"] = list(filtered_df.columns)
                        context["items"] = filtered_df.to_dict('records')
                elif draft.get("items"):
                    # AI가 생성한 표 데이터 사용
                    items_data = draft.get("items", [])
                    if isinstance(items_data, list) and len(items_data) > 0 and isinstance(items_data[0], dict):
                        context["table_headers"] = list(items_data[0].keys())
                        context["items"] = items_data
            except Exception as e:
                st.warning(f"⚠️ 공지문 표 데이터 처리 중 문제: {str(e)}")
            
            template = load_template('gongji_template.html')
            st.session_state[html_key] = generate_html(template, context)
        elif doc_type == '공문':
            draft = gm_data
            context = { "sender_org": draft["sender_org"], "receiver": draft["receiver"], "cc": draft["cc"], "title": draft["title"], "body": text_to_html(draft["body"]), "sender_name": draft["sender_name"], "generation_date": datetime.now().strftime('%Y. %m. %d.') }
            
            # 표 데이터 처리 (AI 생성 또는 사용자 편집)
            try:
                if "df_edited" in gm_data and gm_data["df_edited"] is not None and not gm_data["df_edited"].empty:
                    # 사용자가 편집한 표 데이터 사용
                    filtered_df = gm_data["df_edited"].dropna(how='all')
                    if not filtered_df.empty:
                        context["table_headers"] = list(filtered_df.columns)
                        context["items"] = filtered_df.to_dict('records')
                elif draft.get("items"):
                    # AI가 생성한 표 데이터 사용
                    items_data = draft.get("items", [])
                    if isinstance(items_data, list) and len(items_data) > 0 and isinstance(items_data[0], dict):
                        context["table_headers"] = list(items_data[0].keys())
                        context["items"] = items_data
            except Exception as e:
                st.warning(f"⚠️ 공문 표 데이터 처리 중 문제: {str(e)}")
            
            template = load_template('gongmun_template.html')
            st.session_state[html_key] = generate_html(template, context)
        elif doc_type == '비즈니스 이메일':
            draft = {**e_data, **signature_data}
            context = draft.copy()
            context["signature_company"] = "주식회사 몬쉘코리아"
            
            # 이메일 본문 텍스트 처리 (자연스러운 줄바꿈)
            context["body"] = text_to_html(draft.get("body", ""), for_email=True)
            context["closing"] = text_to_html(draft.get("closing", ""), for_email=True)
            
            # 표 데이터 처리 (AI 생성 또는 사용자 편집)
            try:
                if "df_edited" in e_data and e_data["df_edited"] is not None and not e_data["df_edited"].empty:
                    # 사용자가 편집한 표 데이터 사용
                    filtered_df = e_data["df_edited"].dropna(how='all')
                    if not filtered_df.empty:
                        context["table_headers"] = list(filtered_df.columns)
                        context["items"] = filtered_df.to_dict('records')
                elif e_data.get("items"):
                    # AI가 생성한 표 데이터 사용
                    items_data = e_data.get("items", [])
                    if isinstance(items_data, list) and len(items_data) > 0 and isinstance(items_data[0], dict):
                        context["table_headers"] = list(items_data[0].keys())
                        context["items"] = items_data
            except Exception as e:
                st.warning(f"⚠️ 이메일 표 데이터 처리 중 문제: {str(e)}")
            
            template = load_template('email_template_v2.html')
            st.session_state[html_key] = generate_html(template, context)

if st.session_state.get(html_key):
    st.divider()
    st.subheader("📄 최종 미리보기")
    components.html(st.session_state[html_key], height=600, scrolling=True)
    if doc_type == "비즈니스 이메일":
        st.subheader("📋 복사할 HTML 코드")
        st.code(st.session_state[html_key], language='html')
    else:
        st.divider()
        col1, col2 = st.columns(2)
        with col1:
            pdf_output = generate_pdf(st.session_state[html_key])
            title_for_file = draft.get("title", "document")
            st.download_button(label="📥 PDF 파일로 다운로드", data=pdf_output, file_name=f"{title_for_file}.pdf", mime="application/pdf", use_container_width=True)
        with col2:
            docx_output = generate_docx(draft, doc_type, signature_data)
            st.download_button(label="📄 Word 파일로 다운로드", data=docx_output, file_name=f"{title_for_file}.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)

